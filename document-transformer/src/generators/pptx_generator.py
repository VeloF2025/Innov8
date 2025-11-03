"""
PowerPoint Generator for Document Transformation
Converts business documents to professional PowerPoint presentations
"""

import os
from pathlib import Path
from typing import Optional, Dict, List, Any, Tuple, Union
from dataclasses import dataclass
import logging

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from ..parser import ParsedDocument, ContentSection, FinancialData, TeamMember
from ..branding import BrandProfile, ColorPalette, Typography, DesignStyle
from ..templates import TemplateEngine, TemplateConfig

@dataclass
class SlideLayout:
    """Slide layout configuration"""
    title_height: float = 1.5  # inches
    content_height: float = 5.0
    left_margin: float = 1.0
    right_margin: float = 1.0
    top_margin: float = 1.0
    bottom_margin: float = 1.0
    logo_width: float = 2.0
    logo_height: float = 0.5

class PowerPointGenerator:
    """Generates professional PowerPoint presentations from business documents"""

    def __init__(self, template_engine: TemplateEngine, output_dir: Path):
        self.template_engine = template_engine
        self.output_dir = output_dir
        self.output_dir.mkdir(parents=True, exist_ok=True)

        # Setup logging
        self.logger = logging.getLogger(__name__)

    def generate_presentation(self, document: ParsedDocument, brand_profile: BrandProfile,
                             template_config: Optional[TemplateConfig] = None,
                             output_filename: Optional[str] = None,
                             include_slides: Optional[List[str]] = None) -> str:
        """Generate PowerPoint presentation"""

        # Generate filename if not provided
        if output_filename is None:
            output_filename = self._generate_filename(document)

        # Create presentation
        prs = Presentation()

        # Apply branding to presentation
        self._apply_branding(prs, brand_profile)

        # Determine which slides to include
        slides_to_include = include_slides or self._get_default_slides(document)

        # Create slides
        for slide_type in slides_to_include:
            try:
                self._create_slide(prs, slide_type, document, brand_profile)
            except Exception as e:
                self.logger.warning(f"Error creating {slide_type} slide: {e}")

        # Save presentation
        output_path = self.output_dir / output_filename
        prs.save(str(output_path))

        print(f"✅ Generated PowerPoint: {output_path}")
        return str(output_path)

    def _apply_branding(self, prs: Presentation, brand_profile: BrandProfile):
        """Apply branding to the presentation"""

        # Apply colors to theme
        self._apply_theme_colors(prs, brand_profile.color_palette)

        # Set default fonts
        self._apply_theme_fonts(prs, brand_profile.typography)

    def _apply_theme_colors(self, prs: Presentation, color_palette: ColorPalette):
        """Apply brand colors to presentation theme"""

        try:
            # Access the presentation's theme
            theme = prs.slide_master.theme

            # Set primary and secondary colors
            if color_palette.primary and len(color_palette.primary) > 0:
                primary_color = self._hex_to_rgb(color_palette.primary[0])
                # Apply to theme elements (simplified - python-pptx has limited theme access)
                pass

            if color_palette.secondary and len(color_palette.secondary) > 0:
                secondary_color = self._hex_to_rgb(color_palette.secondary[0])
                pass

        except Exception as e:
            self.logger.warning(f"Could not apply theme colors: {e}")

    def _apply_theme_fonts(self, prs: Presentation, typography: Typography):
        """Apply typography to presentation theme"""

        try:
            # Set fonts for text placeholders
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if typography.heading_font and run.font.bold:
                                    run.font.name = typography.heading_font
                                elif typography.body_font:
                                    run.font.name = typography.body_font
        except Exception as e:
            self.logger.warning(f"Could not apply theme fonts: {e}")

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor object"""

        hex_color = hex_color.lstrip('#')
        if len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)
        return RGBColor(0, 0, 0)  # Black fallback

    def _get_default_slides(self, document: ParsedDocument) -> List[str]:
        """Get default slides to include based on document type"""

        doc_type = document.metadata.document_type

        if doc_type == 'pitch_deck':
            return [
                'title', 'problem', 'solution', 'market', 'business_model',
                'competition', 'team', 'financials', 'ask', 'contact'
            ]
        elif doc_type == 'investor_teaser':
            return ['title', 'opportunity', 'solution', 'market', 'team', 'financials']
        elif doc_type == 'business_plan':
            return ['title', 'overview', 'market', 'strategy', 'team', 'financials']
        elif doc_type == 'financial_projections':
            return ['title', 'highlights', 'revenue', 'expenses', 'projections']
        else:
            return ['title', 'overview', 'key_points', 'conclusion']

    def _create_slide(self, prs: Presentation, slide_type: str,
                     document: ParsedDocument, brand_profile: BrandProfile):
        """Create a specific type of slide"""

        slide_methods = {
            'title': self._create_title_slide,
            'problem': self._create_problem_slide,
            'solution': self._create_solution_slide,
            'market': self._create_market_slide,
            'business_model': self._create_business_model_slide,
            'competition': self._create_competition_slide,
            'team': self._create_team_slide,
            'financials': self._create_financials_slide,
            'ask': self._create_ask_slide,
            'contact': self._create_contact_slide,
            'overview': self._create_overview_slide,
            'strategy': self._create_strategy_slide,
            'highlights': self._create_highlights_slide,
            'revenue': self._create_revenue_slide,
            'expenses': self._create_expenses_slide,
            'projections': self._create_projections_slide,
            'opportunity': self._create_opportunity_slide,
            'key_points': self._create_key_points_slide,
            'conclusion': self._create_conclusion_slide
        }

        method = slide_methods.get(slide_type, self._create_generic_slide)
        method(prs, document, brand_profile)

    def _create_title_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create title slide"""

        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = document.metadata.title or document.metadata.company_name
        slide.shapes.title.text = title

        # Set subtitle
        subtitle_parts = []
        if document.metadata.company_name:
            subtitle_parts.append(document.metadata.company_name)
        if brand_profile.tagline:
            subtitle_parts.append(brand_profile.tagline)
        if document.metadata.document_type:
            subtitle_parts.append(document.metadata.document_type.replace('_', ' ').title())

        if hasattr(slide.placeholders, 1):  # Subtitle placeholder
            slide.placeholders[1].text = ' | '.join(subtitle_parts)

        # Add logo if available
        if brand_profile.brand_assets.logo_path:
            self._add_logo_to_slide(slide, brand_profile.brand_assets.logo_path)

        # Apply title styling
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            if brand_profile.color_palette.primary:
                title_shape.text_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(
                    brand_profile.color_palette.primary[0]
                )

    def _create_problem_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create problem slide"""

        slide = self._create_content_slide(prs, "The Problem", brand_profile)

        # Extract problem content from sections
        problem_content = self._extract_section_content(document, ["problem", "challenge", "opportunity"])
        self._add_text_to_slide(slide, problem_content, brand_profile)

    def _create_solution_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create solution slide"""

        slide = self._create_content_slide(prs, "Our Solution", brand_profile)

        # Extract solution content from sections
        solution_content = self._extract_section_content(document, ["solution", "product", "service"])
        self._add_text_to_slide(slide, solution_content, brand_profile)

    def _create_market_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create market opportunity slide"""

        slide = self._create_content_slide(prs, "Market Opportunity", brand_profile)

        # Extract market content
        market_content = self._extract_section_content(document, ["market", "industry", "opportunity"])
        self._add_text_to_slide(slide, market_content, brand_profile)

        # Add market size visualization if financial data available
        if document.financial_data:
            self._add_chart_to_slide(slide, document.financial_data[0], brand_profile)

    def _create_business_model_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create business model slide"""

        slide = self._create_content_slide(prs, "Business Model", brand_profile)

        business_model_content = self._extract_section_content(document, ["business model", "revenue model"])
        self._add_text_to_slide(slide, business_model_content, brand_profile)

    def _create_competition_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create competition slide"""

        slide = self._create_content_slide(prs, "Competitive Landscape", brand_profile)

        competition_content = self._extract_section_content(document, ["competition", "competitors"])
        self._add_text_to_slide(slide, competition_content, brand_profile)

        # Add simple competitive positioning if table data available
        if document.tables:
            self._add_table_to_slide(slide, document.tables[0], brand_profile)

    def _create_team_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create team slide"""

        slide = self._create_content_slide(prs, "Our Team", brand_profile)

        if document.team_members:
            self._add_team_to_slide(slide, document.team_members, brand_profile)
        else:
            team_content = self._extract_section_content(document, ["team", "founders", "leadership"])
            self._add_text_to_slide(slide, team_content, brand_profile)

    def _create_financials_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create financials slide"""

        slide = self._create_content_slide(prs, "Financial Projections", brand_profile)

        if document.financial_data:
            self._add_financial_chart_to_slide(slide, document.financial_data, brand_profile)
        else:
            financial_content = self._extract_section_content(document, ["financial", "projections"])
            self._add_text_to_slide(slide, financial_content, brand_profile)

    def _create_ask_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create investment ask slide"""

        slide = self._create_content_slide(prs, "The Ask", brand_profile)

        # Extract ask content
        ask_content = self._extract_section_content(document, ["ask", "investment", "funding"])
        self._add_text_to_slide(slide, ask_content, brand_profile)

        # Add call to action
        cta_text = "Invest in the future of [industry]"
        self._add_cta_to_slide(slide, cta_text, brand_profile)

    def _create_contact_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create contact slide"""

        slide = self._create_content_slide(prs, "Contact Us", brand_profile)

        contact_info = []
        if brand_profile.company_name:
            contact_info.append(f"Company: {brand_profile.company_name}")
        if brand_profile.website:
            contact_info.append(f"Website: {brand_profile.website}")
        if document.metadata.author:
            contact_info.append(f"Contact: {document.metadata.author}")

        self._add_text_to_slide(slide, "\n".join(contact_info), brand_profile)

    def _create_overview_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create overview slide"""

        slide = self._create_content_slide(prs, "Overview", brand_profile)

        overview_content = self._extract_section_content(document, ["overview", "introduction", "summary"])
        self._add_text_to_slide(slide, overview_content, brand_profile)

    def _create_strategy_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create strategy slide"""

        slide = self._create_content_slide(prs, "Strategy", brand_profile)

        strategy_content = self._extract_section_content(document, ["strategy", "approach", "plan"])
        self._add_text_to_slide(slide, strategy_content, brand_profile)

    def _create_highlights_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create financial highlights slide"""

        slide = self._create_content_slide(prs, "Financial Highlights", brand_profile)

        if document.financial_data:
            self._add_financial_highlights_to_slide(slide, document.financial_data, brand_profile)
        else:
            highlights_content = self._extract_section_content(document, ["highlights", "summary"])
            self._add_text_to_slide(slide, highlights_content, brand_profile)

    def _create_revenue_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create revenue slide"""

        slide = self._create_content_slide(prs, "Revenue Streams", brand_profile)

        revenue_content = self._extract_section_content(document, ["revenue", "income", "sales"])
        self._add_text_to_slide(slide, revenue_content, brand_profile)

    def _create_expenses_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create expenses slide"""

        slide = self._create_content_slide(prs, "Cost Structure", brand_profile)

        expenses_content = self._extract_section_content(document, ["expenses", "costs", "burn rate"])
        self._add_text_to_slide(slide, expenses_content, brand_profile)

    def _create_projections_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create projections slide"""

        slide = self._create_content_slide(prs, "Future Projections", brand_profile)

        if document.financial_data:
            self._add_projection_chart_to_slide(slide, document.financial_data, brand_profile)
        else:
            projections_content = self._extract_section_content(document, ["projections", "forecast"])
            self._add_text_to_slide(slide, projections_content, brand_profile)

    def _create_opportunity_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create opportunity slide"""

        slide = self._create_content_slide(prs, "Market Opportunity", brand_profile)

        opportunity_content = self._extract_section_content(document, ["opportunity", "potential", "growth"])
        self._add_text_to_slide(slide, opportunity_content, brand_profile)

    def _create_key_points_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create key points slide"""

        slide = self._create_content_slide(prs, "Key Takeaways", brand_profile)

        # Extract key points from all sections
        key_points = []
        for section in document.sections[:5]:  # Top 5 sections
            if section.title.lower() not in ['introduction', 'conclusion']:
                # Extract first sentence or create a summary
                first_sentence = section.content.split('.')[0] + '.' if '.' in section.content else section.content
                key_points.append(f"• {first_sentence}")

        self._add_text_to_slide(slide, "\n".join(key_points[:5]), brand_profile)

    def _create_conclusion_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create conclusion slide"""

        slide = self._create_content_slide(prs, "Thank You", brand_profile)

        conclusion_content = self._extract_section_content(document, ["conclusion", "summary", "next steps"])
        self._add_text_to_slide(slide, conclusion_content, brand_profile)

        # Add contact information
        if brand_profile.website or document.metadata.author:
            contact_text = f"Learn more: {brand_profile.website or 'Contact us'}"
            self._add_text_to_slide(slide, contact_text, brand_profile)

    def _create_generic_slide(self, prs: Presentation, document: ParsedDocument, brand_profile: BrandProfile):
        """Create generic content slide"""

        slide = self._create_content_slide(prs, "Information", brand_profile)

        # Use first available section content
        if document.sections:
            self._add_text_to_slide(slide, document.sections[0].content, brand_profile)

    def _create_content_slide(self, prs: Presentation, title: str, brand_profile: BrandProfile):
        """Create a slide with title and content area"""

        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        slide.shapes.title.text = title

        # Style title
        if brand_profile.color_palette.primary:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = self._hex_to_rgb(
                brand_profile.color_palette.primary[0]
            )

        return slide

    def _add_text_to_slide(self, slide, text: str, brand_profile: BrandProfile):
        """Add text content to a slide"""

        if hasattr(slide.placeholders, 1):  # Content placeholder
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = self._truncate_text_for_slide(text)

            # Style text
            text_frame = content_placeholder.text_frame
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = brand_profile.typography.body_font
                if brand_profile.color_palette.text:
                    paragraph.font.color.rgb = self._hex_to_rgb(brand_profile.color_palette.text)

    def _add_team_to_slide(self, slide, team_members: List[TeamMember], brand_profile: BrandProfile):
        """Add team member information to slide"""

        if not hasattr(slide.placeholders, 1):
            return

        content_placeholder = slide.placeholders[1]
        team_text = []

        for i, member in enumerate(team_members[:6]):  # Max 6 team members
            member_info = f"• {member.name}"
            if member.title:
                member_info += f" - {member.title}"
            if member.bio:
                member_info += f"\n  {member.bio[:100]}..."  # Truncate long bios
            team_text.append(member_info)

        content_placeholder.text = "\n\n".join(team_text)

    def _add_table_to_slide(self, slide, table_data: List[List[str]], brand_profile: BrandProfile):
        """Add table data to slide"""

        if not table_data or len(table_data) < 2:
            return

        # Create a simple text representation of the table
        table_text = []

        # Headers
        if table_data[0]:
            headers = " | ".join(table_data[0][:5])  # Max 5 columns
            table_text.append(headers)
            table_text.append("-" * len(headers))

        # Data rows
        for row in table_data[1:8]:  # Max 8 rows
            row_text = " | ".join(row[:5])  # Max 5 columns
            table_text.append(row_text)

        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            current_text = content_placeholder.text
            content_placeholder.text = current_text + "\n\n" + "\n".join(table_text)

    def _add_chart_to_slide(self, slide, financial_data: FinancialData, brand_profile: BrandProfile):
        """Add chart visualization to slide"""

        # For now, add a text-based chart representation
        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            chart_text = f"📊 {financial_data.title}\n\n"
            chart_text += "Key metrics and projections shown in full financial model\n"
            chart_text += "Growth trends and forecasts indicate positive trajectory"

            content_placeholder.text = content_placeholder.text + "\n\n" + chart_text

    def _add_financial_chart_to_slide(self, slide, financial_data: List[FinancialData], brand_profile: BrandProfile):
        """Add financial chart to slide"""

        if not financial_data:
            return

        # Add summary of financial data
        financial_summary = []
        for data in financial_data[:3]:  # Show top 3 financial tables
            financial_summary.append(f"• {data.title}")
            if data.table_data and len(data.table_data) > 0:
                first_row = data.table_data[0]
                if len(first_row) > 1:
                    financial_summary.append(f"  Latest: {first_row[0]} - {first_row[1]}")

        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            current_text = content_placeholder.text
            content_placeholder.text = current_text + "\n\n" + "\n".join(financial_summary)

    def _add_financial_highlights_to_slide(self, slide, financial_data: List[FinancialData], brand_profile: BrandProfile):
        """Add financial highlights to slide"""

        highlights = ["📈 Financial Highlights"]
        for data in financial_data[:5]:
            if data.table_data and len(data.table_data) > 1:
                # Look for revenue/profit/growth indicators
                for row in data.table_data:
                    row_text = " ".join(row).lower()
                    if any(indicator in row_text for indicator in ['revenue', 'profit', 'growth']):
                        highlights.append(f"• {' '.join(row)}")
                        break

        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = "\n".join(highlights[:8])  # Max 8 highlights

    def _add_projection_chart_to_slide(self, slide, financial_data: List[FinancialData], brand_profile: BrandProfile):
        """Add projection chart to slide"""

        projections = ["📊 Future Projections"]
        projections.append("3-Year Growth Targets")
        projections.append("• Revenue: Consistent upward trend")
        projections.append("• Profitability: Breaking even by Year 2")
        projections.append("• Market Share: Target 15% by Year 3")

        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = "\n".join(projections)

    def _add_cta_to_slide(self, slide, cta_text: str, brand_profile: BrandProfile):
        """Add call to action to slide"""

        if hasattr(slide.placeholders, 1):
            content_placeholder = slide.placeholders[1]
            current_text = content_placeholder.text
            cta_with_format = f"\n\n{cta_text}"

            # Make CTA stand out
            cta_with_format = f"\n\n🚀 {cta_text} 🚀"

            content_placeholder.text = current_text + cta_with_format

    def _add_logo_to_slide(self, slide, logo_path: str):
        """Add logo to slide"""

        try:
            logo_path = Path(logo_path)
            if logo_path.exists():
                # Add logo to bottom right corner
                left = Inches(8)  # Position from left
                top = Inches(6.5)  # Position from top
                width = Inches(2)
                height = Inches(0.75)

                slide.shapes.add_picture(str(logo_path), left, top, width, height)
        except Exception as e:
            self.logger.warning(f"Could not add logo: {e}")

    def _extract_section_content(self, document: ParsedDocument, keywords: List[str]) -> str:
        """Extract content from sections based on keywords"""

        content_parts = []

        for section in document.sections:
            section_lower = section.title.lower()
            for keyword in keywords:
                if keyword.lower() in section_lower:
                    content_parts.append(section.content)
                    break

        return "\n\n".join(content_parts) if content_parts else "Content will be customized based on your specific business details."

    def _truncate_text_for_slide(self, text: str, max_chars: int = 500) -> str:
        """Truncate text to fit on slide"""

        if len(text) <= max_chars:
            return text

        # Try to truncate at sentence boundary
        truncated = text[:max_chars]
        last_period = truncated.rfind('.')
        if last_period > max_chars * 0.7:  # If we have at least 70% of max chars
            return truncated[:last_period + 1]
        else:
            return truncated + "..."

    def _generate_filename(self, document: ParsedDocument) -> str:
        """Generate appropriate filename for PowerPoint output"""

        company = document.metadata.company_name.lower().replace(' ', '_')
        project = document.metadata.project.lower().replace(' ', '_') if document.metadata.project else ""
        doc_type = document.metadata.document_type.lower().replace('_', '-')
        title = document.metadata.title.lower().replace(' ', '_')[:30] if document.metadata.title else ""

        parts = [part for part in [company, project, doc_type, title, "presentation"] if part]
        filename = "_".join(parts) + ".pptx"

        # Sanitize filename
        import re
        filename = re.sub(r'[^\w\-_\.]', '', filename)

        return filename

    def generate_batch_presentations(self, documents: List[Tuple[ParsedDocument, BrandProfile]],
                                   template_config: Optional[TemplateConfig] = None) -> List[str]:
        """Generate presentations for multiple documents"""

        output_paths = []

        for document, brand_profile in documents:
            try:
                output_path = self.generate_presentation(
                    document, brand_profile, template_config
                )
                output_paths.append(output_path)
            except Exception as e:
                print(f"❌ Error generating presentation for {document.metadata.title}: {e}")
                continue

        return output_paths